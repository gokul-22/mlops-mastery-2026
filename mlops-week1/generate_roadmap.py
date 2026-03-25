from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

def create_roadmap_pptx():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MLOps Mastery 2026: Project Roadmap"
    subtitle.text = "High-Performance Image Classification Microservice\nBuilding 'Compound AI Systems' from Scratch"

    # --- Slide 2: What We've Built (Phase 1) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Completed: Engineering Foundation"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Architecture & Environment Setup"
    p = tf.add_paragraph()
    p.text = "• Environment: Python 3.11.15 (Conda) + Poetry"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Config: Type-safe Pydantic V2 Settings"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Model: PyTorch MobileNetV3 with Evaluation Mode"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• API: FastAPI with Schema Validation"
    p.level = 1

    # --- Slide 3: The 6-Month Roadmap ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Journey Ahead: Phases 2-3"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Scaling and Specialization"
    p = tf.add_paragraph()
    p.text = "• Phase 2: Kubernetes, MLflow, and Data Versioning (DVC)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Phase 3: LLMOps, vLLM Serving, and Vector Databases"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Capstone: Autonomous SQL/API Agents"
    p.level = 1

    # --- Slide 4: Key Technologies Used ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Modern MLOps Stack (2026)"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Industry-Standard Tooling"
    p = tf.add_paragraph()
    p.text = "• Core: Python 3.11 / FastAPI / Pydantic"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Infrastructure: Docker / Kubernetes / Helm"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• AI Ops: MLflow / DVC / Prefect"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• LLM: Qdrant / vLLM / LangChain"
    p.level = 1

    # --- Slide 5: The "Self-Healing" System ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Goal: The Self-Healing System"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Autonomous Monitoring & Retraining"
    p = tf.add_paragraph()
    p.text = "• Real-time Drift Detection using Evidently AI"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Automated Retraining triggered by Performance Drops"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Full Traceability: From Data Origin to Model Decision"
    p.level = 1

    filename = 'MLOps_Mastery_Roadmap.pptx'
    prs.save(filename)
    print(f"Successfully generated {filename}")

if __name__ == "__main__":
    create_roadmap_pptx()
